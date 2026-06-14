"""
reports/ppt_generator.py
Professional PowerPoint report — 9 slides with real charts.
Built by Shweta Mishra.
"""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass

import pandas as pd

from analytics.stats_engine import (
    CorrelationResult,
    KPIResult,
    OutlierResult,
    compute_all_descriptive_stats,
)
from core.validator import DataProfile

logger = logging.getLogger(__name__)


@dataclass
class ReportInput:
    df: pd.DataFrame
    profile: DataProfile
    kpi: KPIResult | None
    outliers: dict[str, OutlierResult]
    correlation: CorrelationResult | None
    ai_insights: dict | None
    file_name: str
    metric_col: str


def _make_dist_chart(df: pd.DataFrame, col: str) -> bytes | None:
    """Distribution histogram as PNG bytes."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(8, 4), facecolor="#1E293B")
        ax.set_facecolor("#0F172A")
        data = pd.to_numeric(df[col], errors="coerce").dropna()
        ax.hist(data, bins=30, color="#0D9488", edgecolor="#0F172A", alpha=0.85)
        ax.set_title(f"Distribution of {col}", color="white", fontsize=13, pad=10)
        ax.set_xlabel(col, color="#94A3B8", fontsize=10)
        ax.set_ylabel("Count", color="#94A3B8", fontsize=10)
        ax.tick_params(colors="#94A3B8")
        for spine in ax.spines.values():
            spine.set_edgecolor("#334155")
        ax.axvline(data.mean(), color="#F59E0B", linestyle="--", linewidth=1.5,
                   label=f"Mean: {data.mean():.2f}")
        ax.axvline(data.median(), color="#10B981", linestyle="--", linewidth=1.5,
                   label=f"Median: {data.median():.2f}")
        ax.legend(facecolor="#1E293B", edgecolor="#334155", labelcolor="white",
                  fontsize=9)
        buf = io.BytesIO()
        plt.savefig(buf, format="png", dpi=120, bbox_inches="tight",
                    facecolor="#1E293B")
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception:
        logger.warning("Dist chart failed for col: %s", col)
        return None


def _make_bar_chart(df: pd.DataFrame, cat_col: str, num_col: str) -> bytes | None:
    """Top categories bar chart."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        grouped = (
            df.groupby(cat_col)[num_col]
            .sum()
            .sort_values(ascending=False)
            .head(8)
        )
        if len(grouped) == 0:
            return None

        fig, ax = plt.subplots(figsize=(8, 4), facecolor="#1E293B")
        ax.set_facecolor("#0F172A")
        colors = ["#0D9488", "#14B8A6", "#2DD4BF", "#5EEAD4",
                  "#99F6E4", "#0891B2", "#06B6D4", "#38BDF8"]
        bars = ax.barh(
            grouped.index.astype(str),
            grouped.values,
            color=colors[: len(grouped)],
            edgecolor="#0F172A",
        )
        ax.set_title(f"{num_col} by {cat_col}", color="white", fontsize=13, pad=10)
        ax.tick_params(colors="#94A3B8")
        for spine in ax.spines.values():
            spine.set_edgecolor("#334155")
        ax.set_xlabel(num_col, color="#94A3B8", fontsize=10)
        for bar in bars:
            width = bar.get_width()
            ax.text(
                width * 1.01, bar.get_y() + bar.get_height() / 2,
                f"{width:,.0f}", va="center", color="white", fontsize=8,
            )
        buf = io.BytesIO()
        plt.savefig(buf, format="png", dpi=120, bbox_inches="tight",
                    facecolor="#1E293B")
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception:
        logger.warning("Bar chart failed: %s / %s", cat_col, num_col)
        return None


def _make_corr_heatmap(matrix: pd.DataFrame) -> bytes | None:
    """Correlation heatmap."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(7, 5), facecolor="#1E293B")
        ax.set_facecolor("#1E293B")
        cols = matrix.columns.tolist()
        vals = matrix.values
        im = ax.imshow(vals, cmap="RdBu_r", vmin=-1, vmax=1, aspect="auto")
        ax.set_xticks(range(len(cols)))
        ax.set_yticks(range(len(cols)))
        ax.set_xticklabels(cols, rotation=45, ha="right", color="#94A3B8", fontsize=9)
        ax.set_yticklabels(cols, color="#94A3B8", fontsize=9)
        ax.set_title("Correlation Matrix", color="white", fontsize=13, pad=10)
        for i in range(len(cols)):
            for j in range(len(cols)):
                val = vals[i, j]
                color = "white" if abs(val) > 0.5 else "#94A3B8"
                ax.text(j, i, f"{val:.2f}", ha="center", va="center",
                        color=color, fontsize=8, fontweight="bold")
        plt.colorbar(im, ax=ax).ax.yaxis.set_tick_params(color="#94A3B8")
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format="png", dpi=120, bbox_inches="tight",
                    facecolor="#1E293B")
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception:
        logger.warning("Heatmap failed")
        return None


def generate_ppt_report(report_input: ReportInput) -> bytes:
    """Generate a 9-slide PowerPoint report. Returns bytes."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise ImportError("pip install python-pptx") from e

    ri = report_input
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Colors — lowercase to satisfy ruff N806
    navy   = RGBColor(0x0F, 0x17, 0x2A)
    teal   = RGBColor(0x0D, 0x94, 0x88)
    teal_l = RGBColor(0x14, 0xB8, 0xA6)
    white  = RGBColor(0xF8, 0xFA, 0xFC)
    slate  = RGBColor(0x64, 0x74, 0x8B)
    dark   = RGBColor(0x1E, 0x29, 0x3B)
    red    = RGBColor(0xDC, 0x26, 0x26)
    amber  = RGBColor(0xD9, 0x77, 0x06)
    green  = RGBColor(0x16, 0xA3, 0x4A)

    def rect(slide, left, top, width, height, fill):
        sh = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.fill.background()
        return sh

    def txt(slide, text, left, top, width, height,
            size=12, bold=False, color=None, align="left", italic=False):
        from pptx.enum.text import PP_ALIGN
        color = color or white
        bx = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = bx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return bx

    def add_img(slide, png_bytes, left, top, width, height):
        if png_bytes is None:
            return
        slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            Inches(left), Inches(top), Inches(width), Inches(height),
        )

    branding = "   Built by Shweta Mishra  ·  Excel Auto-Analyst"

    # ── Slide 1: Cover ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, navy)
    rect(s, 0, 0, 0.18, 7.5, teal)
    rect(s, 0, 0, 13.33, 1.4, RGBColor(0x0B, 0x1A, 0x22))
    txt(s, "DATA ANALYSIS REPORT", 0.4, 0.3, 12, 0.6,
        size=11, bold=True, color=teal_l)
    title = (ri.file_name.replace(".csv", "").replace(".xlsx", "")
             .replace("_", " ").title())
    txt(s, title, 0.4, 1.6, 12, 1.3, size=44, bold=True, color=white)
    txt(s,
        f"{ri.profile.total_rows:,} rows  |  "
        f"{ri.profile.total_columns} columns  |  "
        f"Quality: {ri.profile.quality_score}/100",
        0.4, 3.1, 12, 0.5, size=14, color=slate)
    qs = ri.profile.quality_score
    qc = green if qs >= 75 else amber if qs >= 50 else red
    rect(s, 0.4, 3.8, 2.2, 0.75, qc)
    txt(s, f"Quality: {qs}/100", 0.4, 3.8, 2.2, 0.75,
        size=14, bold=True, color=white, align="center")
    rect(s, 3.0, 3.8, 2.5, 0.75, dark)
    txt(s, f"Numeric: {len(ri.profile.numeric_columns)} cols",
        3.0, 3.8, 2.5, 0.75, size=13, color=teal_l, align="center")
    rect(s, 5.8, 3.8, 2.5, 0.75, dark)
    txt(s, f"Categorical: {len(ri.profile.categorical_columns)} cols",
        5.8, 3.8, 2.5, 0.75, size=13, color=teal_l, align="center")
    rect(s, 0, 6.8, 13.33, 0.7, RGBColor(0x06, 0x14, 0x20))
    txt(s, f"  {branding}", 0, 6.85, 13.33, 0.5,
        size=11, color=slate, align="center")

    # ── Slide 2: Data Overview ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, navy)
    rect(s, 0, 0, 13.33, 1.0, dark)
    rect(s, 0, 0, 0.18, 7.5, teal)
    txt(s, "DATA OVERVIEW", 0.4, 0.18, 9, 0.65, size=22, bold=True, color=white)
    txt(s, f"File: {ri.file_name}", 0.4, 0.65, 9, 0.35, size=11, color=slate)
    metrics = [
        ("Total Rows",    f"{ri.profile.total_rows:,}",         teal),
        ("Total Columns", str(ri.profile.total_columns),        teal),
        ("Missing %",     f"{ri.profile.total_missing_pct:.1f}%",
         red if ri.profile.total_missing_pct > 10 else green),
        ("Duplicates",    str(ri.profile.duplicate_row_count),
         red if ri.profile.duplicate_row_count > 0 else green),
        ("Memory (MB)",   str(ri.profile.memory_mb),            teal),
        ("Quality Score", f"{ri.profile.quality_score}/100",
         green if ri.profile.quality_score >= 75 else amber),
    ]
    for i, (label, val, col) in enumerate(metrics):
        bx = 0.4 + (i % 3) * 4.3
        by = 1.2 + (i // 3) * 1.85
        rect(s, bx, by, 3.9, 1.65, dark)
        rect(s, bx, by, 3.9, 0.07, col)
        txt(s, val,   bx, by + 0.2,  3.9, 0.85, size=30, bold=True,
            color=col, align="center")
        txt(s, label, bx, by + 1.1,  3.9, 0.45, size=12,
            color=slate, align="center")
    rect(s, 0.4, 5.1, 12.5, 0.55, RGBColor(0x0B, 0x1A, 0x22))
    txt(s,
        f"Numeric: {ri.profile.numeric_columns[:5]}  |  "
        f"Categorical: {ri.profile.categorical_columns[:5]}",
        0.5, 5.15, 12.3, 0.45, size=10, color=slate)
    txt(s, branding, 0, 7.1, 13.33, 0.35,
        size=9, color=RGBColor(0x33, 0x41, 0x55), align="center")

    # ── Slide 3: KPI Summary ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, navy)
    rect(s, 0, 0, 13.33, 1.0, RGBColor(0x06, 0x3E, 0x31))
    rect(s, 0, 0, 0.18, 7.5, teal)
    txt(s, f"KPI SUMMARY  —  {ri.metric_col.upper()}", 0.4, 0.18, 12, 0.65,
        size=22, bold=True, color=white)
    if ri.kpi:
        kpi_items = [
            ("Total",   f"{ri.kpi.total:,.2f}",   teal),
            ("Mean",    f"{ri.kpi.mean:,.2f}",    teal_l),
            ("Median",  f"{ri.kpi.median:,.2f}",  teal_l),
            ("Std Dev", f"{ri.kpi.std:,.2f}",     slate),
            ("Max",     f"{ri.kpi.max_val:,.2f}",  green),
            ("Min",     f"{ri.kpi.min_val:,.2f}",  red),
            ("Range",   f"{ri.kpi.range_val:,.2f}", amber),
            ("CV %",    f"{ri.kpi.cv_pct:.1f}%",  slate),
        ]
        for i, (label, val, col) in enumerate(kpi_items):
            bx = 0.4 + (i % 4) * 3.2
            by = 1.2 + (i // 4) * 1.9
            rect(s, bx, by, 2.9, 1.7, dark)
            rect(s, bx, by, 2.9, 0.07, col)
            txt(s, val,   bx, by + 0.2, 2.9, 0.9,  size=26, bold=True,
                color=col, align="center")
            txt(s, label, bx, by + 1.2, 2.9, 0.4,  size=12,
                color=slate, align="center")
        if ri.kpi.mom_change_pct is not None:
            mom_c = green if ri.kpi.mom_change_pct >= 0 else red
            sign  = "+" if ri.kpi.mom_change_pct >= 0 else ""
            rect(s, 0.4, 5.2, 4.5, 0.65, mom_c)
            txt(s, f"Month-over-Month: {sign}{ri.kpi.mom_change_pct:.1f}%",
                0.4, 5.2, 4.5, 0.65, size=14, bold=True,
                color=white, align="center")
    else:
        txt(s, "No numeric data available.", 0.4, 2.0, 12, 0.5,
            size=14, color=slate)
    txt(s, branding, 0, 7.1, 13.33, 0.35,
        size=9, color=RGBColor(0x33, 0x41, 0x55), align="center")

    # ── Slide 4: Distribution Chart ───────────────────────────────
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, navy)
    rect(s, 0, 0, 13.33, 1.0, dark)
    rect(s, 0, 0, 0.18, 7.5, teal)
    txt(s, f"DISTRIBUTION  —  {ri.metric_col.upper()}", 0.4, 0.18, 12, 0.65,
        size=20, bold=True, color=white)
    desc_stats = compute_all_descriptive_stats(ri.df, ri.profile)
    if ri.metric_col in desc_stats:
        ds = desc_stats[ri.metric_col]
        stat_rows = [
            ("Mean",    f"{ds.mean:,.2f}"),
            ("Median",  f"{ds.median:,.2f}"),
            ("Std Dev", f"{ds.std:,.2f}"),
            ("Skewness",f"{ds.skewness:,.3f}"),
            ("P25",     f"{ds.p25:,.2f}"),
            ("P75",     f"{ds.p75:,.2f}"),
            ("IQR",     f"{ds.iqr:,.2f}"),
            ("Normal?", "Yes" if ds.is_normal else "No"),
        ]
        rect(s, 0.4, 1.1, 3.2, 5.7, dark)
        txt(s, "Statistics", 0.5, 1.15, 3.0, 0.45,
            size=13, bold=True, color=teal_l)
        for i, (label, val) in enumerate(stat_rows):
            by = 1.65 + i * 0.6
            txt(s, label, 0.5,  by, 1.4, 0.5, size=10, color=slate)
            txt(s, val,   1.9,  by, 1.6, 0.5, size=10, bold=True, color=white)
    dist_png = (
        _make_dist_chart(ri.df, ri.metric_col) if ri.profile.has_numeric else None
    )
    if dist_png:
        add_img(s, dist_png, 3.8, 1.1, 9.1, 5.0)
    else:
        txt(s, "Chart not available.", 3.8, 3.0, 9.1, 0.5, size=14, color=slate)
    txt(s, branding, 0, 7.1, 13.33, 0.35,
        size=9, color=RGBColor(0x33, 0x41, 0x55), align="center")

    # ── Slide 5: Category Breakdown ───────────────────────────────
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, navy)
    rect(s, 0, 0, 13.33, 1.0, dark)
    rect(s, 0, 0, 0.18, 7.5, teal)
    txt(s, "CATEGORY BREAKDOWN", 0.4, 0.18, 12, 0.65, size=22, bold=True, color=white)
    if ri.profile.has_categorical and ri.profile.has_numeric:
        cat_col = ri.profile.categorical_columns[0]
        bar_png = _make_bar_chart(ri.df, cat_col, ri.metric_col)
        add_img(s, bar_png, 0.4, 1.1, 9.0, 5.0)
        rect(s, 9.6, 1.1, 3.4, 5.0, dark)
        txt(s, f"Top {cat_col}", 9.7, 1.15, 3.2, 0.4,
            size=12, bold=True, color=teal_l)
        grouped = (
            ri.df.groupby(cat_col)[ri.metric_col]
            .sum()
            .sort_values(ascending=False)
            .head(7)
        )
        for i, (label, val) in enumerate(grouped.items()):
            by = 1.65 + i * 0.6
            txt(s, str(label)[:15], 9.7,  by, 1.8, 0.5, size=10, color=white)
            txt(s, f"{val:,.0f}",   11.5, by, 1.4, 0.5, size=10,
                color=teal_l, align="right")
    else:
        txt(s, "Requires both numeric and categorical columns.",
            0.4, 2.5, 12, 0.5, size=14, color=slate)
    txt(s, branding, 0, 7.1, 13.33, 0.35,
        size=9, color=RGBColor(0x33, 0x41, 0x55), align="center")

    # ── Slide 6: Outlier Analysis ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, navy)
    rect(s, 0, 0, 13.33, 1.0, RGBColor(0x45, 0x0A, 0x0A))
    rect(s, 0, 0, 0.18, 7.5, red)
    txt(s, "OUTLIER ANALYSIS", 0.4, 0.18, 12, 0.65, size=22, bold=True, color=white)
    txt(s, "IQR method — flagged, not automatically removed",
        0.4, 0.65, 12, 0.35, size=11, color=slate)
    if ri.outliers:
        rect(s, 0.4, 1.1, 12.5, 0.42, dark)
        headers  = ["Column", "Outliers", "Outlier %", "Lower Fence",
                    "Upper Fence", "Severity"]
        x_pos    = [0.5, 3.2, 5.5, 7.2, 9.5, 11.2]
        for header, xp in zip(headers, x_pos, strict=False):
            txt(s, header, xp, 1.1, 1.8, 0.42,
                size=10, bold=True, color=teal_l)
        for i, (col, out) in enumerate(list(ri.outliers.items())[:8]):
            by     = 1.6 + i * 0.58
            row_bg = dark if i % 2 == 0 else RGBColor(0x16, 0x20, 0x30)
            rect(s, 0.4, by, 12.5, 0.52, row_bg)
            sev_col = (
                red if out.outlier_pct > 5 else amber if out.outlier_pct > 1 else green
            )
            lb  = f"{out.lower_bound:.2f}" if out.lower_bound is not None else "-"
            ub  = f"{out.upper_bound:.2f}" if out.upper_bound is not None else "-"
            sev = (
                "High"
                if out.outlier_pct > 5
                else "Medium"
                if out.outlier_pct > 1
                else "Low"
            )
            vals = [col[:18], str(out.outlier_count),
                    f"{out.outlier_pct:.1f}%", lb, ub, sev]
            for j, (val, xp) in enumerate(zip(vals, x_pos, strict=False)):
                c = sev_col if j == 5 else white
                txt(s, val, xp, by + 0.06, 1.8, 0.4, size=9, color=c)
    else:
        txt(s, "No outliers detected.", 0.4, 2.5, 12, 0.5, size=14, color=slate)
    txt(s, branding, 0, 7.1, 13.33, 0.35,
        size=9, color=RGBColor(0x33, 0x41, 0x55), align="center")

    # ── Slide 7: Correlation Heatmap ──────────────────────────────
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, navy)
    rect(s, 0, 0, 13.33, 1.0, RGBColor(0x1E, 0x1B, 0x4B))
    rect(s, 0, 0, 0.18, 7.5, RGBColor(0x60, 0x5C, 0xF6))
    txt(s, "CORRELATION ANALYSIS", 0.4, 0.18, 12, 0.65, size=22, bold=True, color=white)
    if ri.correlation and len(ri.profile.numeric_columns) >= 2:
        corr_png = _make_corr_heatmap(ri.correlation.matrix)
        add_img(s, corr_png, 0.4, 1.1, 7.5, 5.5)
        rect(s, 8.1, 1.1, 4.8, 5.5, dark)
        txt(s, "Strong Correlations", 8.2, 1.15, 4.6, 0.45,
            size=13, bold=True, color=teal_l)
        if ri.correlation.strong_pairs:
            for i, (col_a, col_b, val) in enumerate(ri.correlation.strong_pairs[:7]):
                by    = 1.7 + i * 0.7
                col_c = red if val < -0.5 else green if val > 0.5 else amber
                icon  = "Strong" if abs(val) > 0.8 else "Moderate"
                txt(s, f"{icon}: {col_a[:8]} x {col_b[:8]}",
                    8.2, by, 3.5, 0.35, size=10, color=white)
                txt(s, f"r = {val:.3f}", 8.2, by + 0.35, 3.5, 0.3,
                    size=10, bold=True, color=col_c)
        else:
            txt(s, "No strong correlations (|r| < 0.5)",
                8.2, 2.0, 4.5, 1.0, size=11, color=slate)
    else:
        txt(s, "Needs at least 2 numeric columns.", 0.4, 2.5, 12, 0.5,
            size=14, color=slate)
    txt(s, branding, 0, 7.1, 13.33, 0.35,
        size=9, color=RGBColor(0x33, 0x41, 0x55), align="center")

    # ── Slide 8: AI Insights ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, navy)
    rect(s, 0, 0, 13.33, 1.0, RGBColor(0x4A, 0x19, 0x42))
    rect(s, 0, 0, 0.18, 7.5, RGBColor(0xA8, 0x55, 0xF7))
    txt(
        s, "AI-GENERATED INSIGHTS", 0.4, 0.18, 12, 0.65, size=22, bold=True, color=white
    )
    txt(s, "Powered by Groq LLaMA-3.3", 0.4, 0.65, 12, 0.35, size=11, color=slate)
    if ri.ai_insights:
        headline = ri.ai_insights.get("headline", "")
        rect(s, 0.4, 1.1, 12.5, 0.75, RGBColor(0x2D, 0x1B, 0x45))
        txt(s, headline, 0.5, 1.12, 12.3, 0.65, size=14,
            bold=True, color=RGBColor(0xE9, 0xD5, 0xFF))
        for i, insight in enumerate(ri.ai_insights.get("insights", [])[:3]):
            bx = 0.4 + i * 4.32
            rect(s, bx, 2.05, 4.05, 3.5, dark)
            rect(s, bx, 2.05, 4.05, 0.08, RGBColor(0xA8, 0x55, 0xF7))
            txt(s, insight.get("title", ""),  bx + 0.12, 2.2,  3.8, 0.45,
                size=12, bold=True, color=RGBColor(0xC4, 0xB5, 0xFD))
            txt(s, insight.get("detail", ""), bx + 0.12, 2.73, 3.8, 2.7,
                size=10, color=white)
        rec = ri.ai_insights.get("recommendation", "")
        if rec:
            rect(s, 0.4, 5.75, 12.5, 0.75, green)
            txt(s, f"Recommendation: {rec}", 0.55, 5.78, 12.3, 0.65,
                size=12, bold=True, color=white)
    else:
        txt(s, "AI insights not generated — add GROQ_API_KEY to enable.",
            0.4, 2.0, 12, 0.5, size=14, color=slate, italic=True)
    txt(s, branding, 0, 7.1, 13.33, 0.35,
        size=9, color=RGBColor(0x33, 0x41, 0x55), align="center")

    # ── Slide 9: Closing ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.33, 7.5, navy)
    rect(s, 0, 0, 0.18, 7.5, teal)
    rect(s, 0, 0, 13.33, 0.08, teal)
    txt(s, "Thank You", 0.4, 1.2, 12, 1.2, size=52, bold=True,
        color=white, align="center")
    txt(s, "Report generated by Excel Auto-Analyst",
        0.4, 2.6, 12, 0.5, size=16, color=slate, align="center")
    txt(s,
        f"Dataset: {ri.file_name}  |  "
        f"Rows: {ri.profile.total_rows:,}  |  "
        f"Quality: {ri.profile.quality_score}/100",
        0.4, 3.3, 12, 0.5, size=12, color=slate, align="center")
    rect(s, 3.5, 4.5, 6.3, 1.0, dark)
    txt(s, "★ ★ ★ ★ ★", 3.5, 4.55, 6.3, 0.45, size=22,
        color=amber, align="center")
    txt(s, "Built by Shweta Mishra", 3.5, 5.0, 6.3, 0.45,
        size=13, bold=True, color=teal_l, align="center")
    txt(s, "github.com/Shweta-Mishra-ai/excel-auto-analyst",
        0.4, 6.5, 12, 0.4, size=11, color=slate, align="center")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info("PPT generated: %d bytes", buf.getbuffer().nbytes)
    return buf.read()
